import os
import re
import io
import time
import streamlit as st
import PyPDF2
import docx
from pptx import Presentation
import trafilatura
import requests
from dotenv import load_dotenv
from google import genai
from google.genai import types
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
import markdown
from xhtml2pdf import pisa

# Force reload of environment variables
load_dotenv(override=True)

# Initialize Gemini Client
client = genai.Client(api_key=os.environ.get("GEMINI_API_KEY"))

# ==========================================
# 1. THE IN-MEMORY OMNI-READER (DRIVE API)
# ==========================================
def fetch_dynamic_knowledge(folder_url):
    match = re.search(r'/folders/([a-zA-Z0-9_-]+)', folder_url)
    folder_id = match.group(1) if match else None
    
    if not folder_id:
        return "Error: Invalid Google Drive Link."

    api_key = os.environ.get("GOOGLE_DRIVE_API_KEY")
    service = build('drive', 'v3', developerKey=api_key)
    
    results = service.files().list(q=f"'{folder_id}' in parents and trashed=false").execute()
    items = results.get('files', [])
    
    knowledge_text = "### COMPANY KNOWLEDGE BASE ###\n"
    
    progress_bar = st.progress(0)
    status_text = st.empty()
    total_files = len(items)
    
    if total_files == 0:
        return "Error: No supported files found in the folder."

    for i, item in enumerate(items):
        mime_type = item['mimeType']
        
        if mime_type in [
            'application/pdf', 
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        ]:
            status_text.text(f"Downloading & Parsing: {item['name']}...")
                
            try:
                request = service.files().get_media(fileId=item['id'])
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    status, done = downloader.next_chunk()
                
                fh.seek(0)
                
                if mime_type == 'application/pdf':
                    reader = PyPDF2.PdfReader(fh)
                    for page in reader.pages:
                        ext = page.extract_text()
                        if ext: knowledge_text += ext + "\n"
                
                elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    doc = docx.Document(fh)
                    for para in doc.paragraphs:
                        if para.text.strip(): knowledge_text += para.text + "\n"
                        
                elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    prs = Presentation(fh)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                knowledge_text += shape.text + "\n"
                                
                time.sleep(2.0) # Anti-bot pause
                
            except Exception as e:
                print(f"Skipping {item['name']} due to error: {e}")
                continue 
                
        progress_bar.progress((i + 1) / total_files)
            
    status_text.text("Knowledge Base successfully loaded into memory!")
    time.sleep(1)
    progress_bar.empty()
    status_text.empty()
                
    return knowledge_text

# ==========================================
# 2. AI ENGINE: SUMMARIZER & DYNAMIC SOLVER
# ==========================================
def generate_kb_summary(kb_text):
    """Generates a brief summary of the ingested Knowledge Base."""
    prompt = f"""
    Analyze the following company knowledge base. 
    Write a concise, 100-word summary explaining what this company does, its core products/services, and its primary value proposition.
    
    Knowledge Base:
    {kb_text}
    """
    response = client.models.generate_content(model="gemini-2.5-pro", contents=prompt)
    return response.text

def dynamic_problem_solver(kb_text, problem_statement):
    """
    STRICT BUSINESS LOGIC:
    Checks if the prompt is business-related.
    If YES -> Executes with Web Search + Knowledge Base.
    If NO -> Rejects the request with a specific protocol violation message.
    """
    prompt = f"""
    You are the Conversely Autonomous Business Engine.
    Your existence is dedicated SOLELY to corporate strategy, B2B sales, tender acquisition, and business problem solving.
    
    ### COMPANY KNOWLEDGE BASE:
    {kb_text}
    
    ### USER COMMAND:
    "{problem_statement}"
    
    ### PROTOCOL (STRICT COMPLIANCE REQUIRED):
    
    STEP 1: CLASSIFY THE INTENT
    - Is this request related to business, strategy, sales, coding for business automation, tenders, market research, or the company's data?
    - YES: Proceed to Step 2.
    - NO (e.g., "Tell me a joke", "Write a poem", "How is the weather", "Who is the president", casual chat): 
      STOP IMMEDIATELY. Output EXACTLY this message:
      "**PROTOCOL VIOLATION:** I am a specialized Business Strategy Engine designed for corporate execution. I cannot assist with non-business queries."
    
    STEP 2: EXECUTE (Only if Step 1 is YES)
    - Use the 'google_search' tool to find real-time live data (tenders, startups, news) if needed.
    - Synthesize the search results with the COMPANY KNOWLEDGE BASE.
    - Output a professional, operational strategic brief.
    - NO FLUFF. NO "I hope this helps". Just the data and the strategy.
    
    EXECUTE PROTOCOL NOW.
    """
    
    response = client.models.generate_content(
        model="gemini-2.5-pro", 
        contents=prompt,
        config=types.GenerateContentConfig(
            tools=[{"google_search": {}}]
        )
    )
    return response.text

def dispatch_proposal(sender_email, app_password, receiver_email, problem_statement, ai_solution, company_name="Conversely"):
    # 1. Convert AI Markdown output to HTML for the PDF
    # We add extensions to ensure lists and bolding render perfectly
    html_body = markdown.markdown(ai_solution, extensions=['extra', 'nl2br'])
    
    # 2. Build the Dynamic PDF Template (Upgraded Enterprise CSS)
    pdf_html = f"""
    <html>
    <head>
    <style>
        @page {{
            size: A4; 
            margin: 2.5cm; 
            margin-top: 3cm; 
            margin-bottom: 3cm;
            @frame header_frame {{ -pdf-frame-content: header_content; top: 1cm; margin-left: 2.5cm; margin-right: 2.5cm; height: 1.5cm; border-bottom: 1px solid #E2E8F0; }}
            @frame footer_frame {{ -pdf-frame-content: footer_content; bottom: 1cm; margin-left: 2.5cm; margin-right: 2.5cm; height: 1cm; }}
        }}
        body {{ 
            font-family: "Helvetica Neue", Helvetica, Arial, sans-serif; 
            font-size: 10.5pt; 
            line-height: 1.8; 
            color: #334155; 
        }}
        h1, h2, h3 {{ 
            color: #0F172A; 
            margin-top: 25px; 
            margin-bottom: 12px; 
        }}
        h1 {{ font-size: 16pt; border-bottom: 2px solid #3B82F6; padding-bottom: 5px; }}
        h2 {{ font-size: 14pt; color: #1E293B; }}
        h3 {{ font-size: 12pt; color: #0F172A; font-weight: bold; }}
        strong {{ color: #0F172A; font-weight: bold; }}
        p {{ margin-bottom: 14px; text-align: left; }}
        ul {{ margin-bottom: 14px; margin-left: 15px; padding-left: 15px; list-style-type: disc; }}
        li {{ margin-bottom: 8px; line-height: 1.6; }}
        
        /* Premium Header/Footer Styling */
        .header-text {{ text-align: right; color: #64748B; font-size: 8.5pt; font-weight: bold; text-transform: uppercase; letter-spacing: 1.5px; padding-bottom: 10px; }}
        .footer-text {{ text-align: center; font-size: 8pt; color: #94A3B8; border-top: 1px solid #E2E8F0; padding-top: 8px; }}
        
        /* Modern Title Box */
        .title-box {{ background-color: #F8FAFC; padding: 20px; border-left: 5px solid #3B82F6; margin-bottom: 30px; border-radius: 4px; }}
        .title-box h2 {{ margin: 0 0 10px 0; color: #0F172A; border: none; padding: 0; font-size: 18pt; }}
        .title-box p {{ margin: 0; font-size: 10pt; color: #475569; }}
    </style>
    </head>
    <body>
        <div id="header_content" class="header-text">{company_name} Internal Intelligence & Strategy Brief</div>
        <div id="footer_content" class="footer-text">{company_name} Confidential & Proprietary | Page <pdf:pagenumber> of <pdf:pagecount></div>
        
        <div class="title-box">
            <h2>Strategic Proposal & Solution Analysis</h2>
            <p><strong>Objective:</strong> {problem_statement}</p>
        </div>
        
        <div class="content">
            {html_body}
        </div>
    </body>
    </html>
    """

    # 3. Create PDF in Memory
    pdf_buffer = io.BytesIO()
    pisa_status = pisa.CreatePDF(pdf_html, dest=pdf_buffer)
    pdf_buffer.seek(0)
    
    if pisa_status.err:
        return False, "Failed to generate PDF."

    # 4. Construct the Email
    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = receiver_email
    msg['Subject'] = f"New Strategic Proposal Generated: {problem_statement[:40]}..."
    
    # 5. Email Body with DYNAMIC HTML CTA Button
    # Dynamically extract a first name from the email (e.g., nitin.sharma@... becomes "Nitin Sharma")
    receiver_name = receiver_email.split('@')[0].replace('.', ' ').title()

    email_html = f"""
    <div style="font-family: Arial, sans-serif; color: #333; line-height: 1.6; max-width: 600px; margin: 0 auto; padding: 20px; border: 1px solid #E2E8F0; border-radius: 8px;">
        <h2 style="color: #0F172A; border-bottom: 2px solid #3B82F6; padding-bottom: 10px;">New Autonomous SDR Proposal</h2>
        
        <p>Hi {receiver_name},</p>
        <p>The AI Engine has completed its web research and generated a highly targeted strategic proposal based on the {company_name} knowledge base.</p>
        
        <div style="background-color: #F8FAFC; padding: 15px; border-left: 4px solid #3B82F6; margin: 20px 0;">
            <p style="margin: 0;"><strong>Target/Problem:</strong><br> {problem_statement}</p>
        </div>
        
        <p>Please find the fully formatted PDF proposal attached for your review. You can authorize the next phase of outreach directly below.</p>
        
        <div style="text-align: center; margin: 35px 0;">
            <a href="mailto:{sender_email}?subject=Re: Strategy Approval for {problem_statement[:30]}&body=Hi Team,%0D%0A%0D%0AI have reviewed the attached proposal and approve moving forward with this strategy.%0D%0A%0D%0ABest,%0D%0A{receiver_name}" 
               style="background-color: #3B82F6; color: #ffffff; padding: 14px 28px; text-decoration: none; border-radius: 6px; font-weight: bold; display: inline-block; font-size: 16px;">
               Review & Approve Strategy
            </a>
        </div>
        
        <p style="font-size: 0.85em; color: #64748B; text-align: center; border-top: 1px solid #E2E8F0; padding-top: 15px;">
            Automated Dispatch by {company_name} Agentic System ⚡
        </p>
    </div>
    """
    
    msg.attach(MIMEText(email_html, 'html'))
    
    # 6. Attach the Generated PDF
    safe_filename = problem_statement[:20].replace(' ', '_').replace('/', '_')
    pdf_attachment = MIMEApplication(pdf_buffer.read(), _subtype="pdf")
    pdf_attachment.add_header('Content-Disposition', 'attachment', filename=f"{company_name}_Proposal_{safe_filename}.pdf")
    msg.attach(pdf_attachment)

    # 7. Send the Email
    try:
        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()
        server.login(sender_email, app_password)
        server.send_message(msg)
        server.quit()
        return True, "Formatted Email & Premium PDF successfully dispatched!"
    except Exception as e:
        return False, f"Failed to send email: {e}"

# ==========================================
# 3. STREAMLIT UI DASHBOARD (STATEFUL)
# ==========================================
st.set_page_config(page_title="Conversely Agentic System", page_icon="⚡", layout="wide")

# Initialize Session State Variables
if "kb_text" not in st.session_state:
    st.session_state.kb_text = None
if "kb_summary" not in st.session_state:
    st.session_state.kb_summary = None

st.title("Conversely: Agentic Business Engine")
st.markdown("Ingest organizational knowledge, then deploy the AI agent to dynamically solve problems, find tenders, or generate leads over the web.")

# PHASE 1: KNOWLEDGE INGESTION
st.header("Step 1: Ingest Knowledge Base")
col1, col2 = st.columns([3, 1])
with col1:
    drive_link = st.text_input("Company Data Dump (Google Drive Link)", placeholder="https://drive.google.com/drive/folders/...")
with col2:
    st.markdown("<br>", unsafe_allow_html=True) # Alignment fix
    if st.button("Ingest Data", type="primary", use_container_width=True):
        if not drive_link:
            st.error("Please provide a Google Drive link.")
        else:
            with st.spinner("Ingesting Cloud Knowledge Base..."):
                raw_text = fetch_dynamic_knowledge(drive_link)
                if "Error" in raw_text:
                    st.error(raw_text)
                else:
                    st.session_state.kb_text = raw_text
                    with st.spinner("Generating AI Understanding..."):
                        st.session_state.kb_summary = generate_kb_summary(raw_text)
                    st.rerun() # Refresh UI to show Phase 2

# PHASE 2: DYNAMIC PROBLEM SOLVING
if st.session_state.kb_summary:
    st.success("Knowledge Base Active")
    
    with st.expander("What the AI understands about your company:", expanded=True):
        st.info(st.session_state.kb_summary)

    st.markdown("---")
    st.header("Step 2: Deploy the Agent")
    st.markdown("Ask the agent to solve a problem. Examples:\n* *\"Find 3 live government tenders related to Agentic LLMs and write a pitch for them.\"*\n* *\"Who are our top 3 competitors in Omni-channel CRMs right now, and how do we beat them?\"*")
    
    problem_input = st.text_area("What problem do you want to solve today?", height=100)
    
    # 1. Initialize a state variable to hold the AI's answer
    if "ai_solution" not in st.session_state:
        st.session_state.ai_solution = None

    # 2. Generate the solution and save it to state
    if st.button("Execute Agentic Workflow", type="primary"):
        if not problem_input:
            st.warning("Please enter a problem statement.")
        else:
            with st.spinner("Agent is thinking and searching the live web..."):
                # Save to session_state so it doesn't disappear on the next button click!
                st.session_state.ai_solution = dynamic_problem_solver(st.session_state.kb_text, problem_input)

    # 3. If a solution exists in memory, display it AND show the email form independently
    # OUTPUT LOGIC (Updated for Strict Mode)
    if st.session_state.ai_solution:
        
        # CHECK: Did the AI refuse the task?
        if "PROTOCOL VIOLATION" in st.session_state.ai_solution:
            st.error(st.session_state.ai_solution) # Show as Red Error Box
        
        else:
            # If valid business response, show Green Success & Email Options
            st.markdown("### Agent Output")
            st.write(st.session_state.ai_solution) 

            st.markdown("---")
            st.markdown("### Dispatch Proposal")
            st.markdown("Generate a corporate PDF from this solution and email it directly to a stakeholder.")

            col_email, col_send = st.columns([3, 1])

            with col_email:
                target_email = st.text_input("Recipient Email Address", placeholder="nitin@conversely.com")

            with col_send:
                st.markdown("<br>", unsafe_allow_html=True) 
                if st.button("Email PDF Proposal", use_container_width=True):
                    if not target_email:
                        st.warning("Please enter an email address.")
                    else:
                        sender = os.environ.get("SENDER_EMAIL") 
                        pwd = os.environ.get("EMAIL_APP_PASSWORD")
                        
                        if not sender or not pwd:
                            st.error("Missing Email Credentials! Check your .env file.")
                        else:
                            with st.spinner("Compiling PDF and contacting SMTP server..."):
                                success, message = dispatch_proposal(
                                    sender_email=sender,
                                    app_password=pwd,
                                    receiver_email=target_email,
                                    problem_statement=problem_input,
                                    ai_solution=st.session_state.ai_solution,
                                    company_name="Conversely"
                                )
                                
                                if success:
                                    st.success(message)
                                    st.balloons()
                                else:
                                    st.error(message)