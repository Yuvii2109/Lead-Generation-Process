import os
import re
import io
import PyPDF2
from dotenv import load_dotenv
from google import genai
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import trafilatura
import docx
from pptx import Presentation

load_dotenv(override=True)

# ==========================================
# CONFIGURATION
# ==========================================
CONVERSELY_DRIVE_LINK = "https://drive.google.com/drive/folders/1dzYDW-Nidq405mLFMlmj903qIksvCWwy?usp=sharing"

# Let's test it on a well-known modern startup to see how the AI reacts
TEST_STARTUP_URL = "https://linear.app" 

# ==========================================
# 1. READ CONVERSELY DOCS (OMNI-READER)
# ==========================================
def fetch_conversely_knowledge(folder_path):
    print(f"Fetching ALL Conversely Docs from local folder: {folder_path}...")
    knowledge_text = "### CONVERSELY SERVICES & CAPABILITIES ###\n"
    
    if not os.path.exists(folder_path):
        print(f"Error: Could not find the folder '{folder_path}'.")
        return ""

    # Read EVERY supported file in the folder
    for filename in os.listdir(folder_path):
        filepath = os.path.join(folder_path, filename)
        
        try:
            # --- 1. PDF HANDLER ---
            if filename.endswith('.pdf'):
                print(f"   -> Reading PDF: {filename}")
                with open(filepath, 'rb') as fh:
                    reader = PyPDF2.PdfReader(fh)
                    for page in reader.pages:
                        extracted = page.extract_text()
                        if extracted: knowledge_text += extracted + "\n"
                        
            # --- 2. WORD DOC HANDLER ---
            elif filename.endswith('.docx'):
                print(f"   -> Reading DOCX: {filename}")
                doc = docx.Document(filepath)
                for para in doc.paragraphs:
                    if para.text.strip():
                        knowledge_text += para.text + "\n"
                        
            # --- 3. POWERPOINT HANDLER ---
            elif filename.endswith('.pptx'):
                print(f"   -> Reading PPTX: {filename}")
                prs = Presentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        # Only extract text from shapes that actually contain text
                        if hasattr(shape, "text") and shape.text.strip():
                            knowledge_text += shape.text + "\n"
                            
        except Exception as e:
            print(f"   -> Could not read {filename}: {e}")
            
    print("Knowledge base fully loaded!")
    return knowledge_text

# ==========================================
# 2. SCRAPE THE STARTUP
# ==========================================
def scrape_startup_website(url):
    print(f"\nScraping Startup Website: {url}")
    downloaded = trafilatura.fetch_url(url)
    if downloaded:
        text = trafilatura.extract(downloaded)
        print("   -> Successfully extracted core website text!")
        return text
    else:
        print("   -> Failed to scrape website.")
        return None

# ==========================================
# 3. AI PROPOSAL GENERATION
# ==========================================
def generate_conversely_pitch(startup_text, conversely_docs):
    print("\nAI generating Conversely B2B Proposal...")
    client = genai.Client(api_key="")
    
    prompt = f"""
You are an Elite Tech Sales SDR and Solutions Architect for 'Conversely', a Technical Software Services and SaaS company.
Your goal is to read the website text of a newly discovered startup, figure out what they do, and write a highly targeted, cold-outreach business proposal offering Conversely's services to help them scale.

### CONVERSELY'S TECHNICAL CAPABILITIES:
{conversely_docs}

### STARTUP'S WEBSITE TEXT:
{startup_text}

### INSTRUCTIONS:
1. Analyze the startup. What is their core product? What technical bottlenecks (e.g., scaling infrastructure, UI/UX, database management, cloud architecture) are they likely facing as a new company?
2. Map their likely pain points strictly to Conversely's services mentioned in the provided documentation.
3. Write a crisp, highly professional 4-part email proposal:
   - **Subject Line:** Catchy, tech-focused, and personalized.
   - **The Hook:** Acknowledge what they are building and why it's impressive.
   - **The Value Prop (Bullet Points):** Propose 2-3 specific ways Conversely's technical services can accelerate their roadmap. (Use exact terms from the Conversely docs).
   - **Call to Action:** A low-friction request for a quick 10-minute discovery call.

CRITICAL: Do not use generic placeholders like [Startup Name] or [Contact Name]. You must deduce the name of the company from the provided text and seamlessly integrate it into the email.
"""
    response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    return response.text

# ==========================================
# EXECUTE TEST
# ==========================================
if __name__ == "__main__":
    # Point this to the name of the folder on your laptop
    docs = fetch_conversely_knowledge("Conversely_Docs") 
    startup_info = scrape_startup_website(TEST_STARTUP_URL)
    
    if startup_info and docs:
        pitch = generate_conversely_pitch(startup_info, docs)
        print("\n" + "="*60)
        print("CONVERSELY PROPOSAL GENERATED:")
        print("="*60)
        print(pitch)
